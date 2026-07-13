#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
qa2.py — PPT Pro Studio QA 2.1 layout & quality audit.

Goes beyond "did it render?" to "does it look right and is it safe to ship?":
  * bounds      — no shape escapes the slide canvas
  * placeholders — no leftover [..], 图片缺失, TODO, 占位 tokens
  * font floor  — body text >= 11pt; only footer/page numbers may be 9pt
  * transitions — every slide has a <p:transition> (翻页动画)
  * icon/vector — count of native custGeom (embedded vector icons)
  * brand theme — slide-master theme1.xml carries a CUSTOM accent1 colour
  * cjk font    — CJK text must not sit on a Latin-only font (would box/render)
  * contrast    — text vs its box fill must clear a minimal WCAG ratio
  * truncation  — flag runs cut with … / ... (content was dropped)
  * image scale — warn when an embedded picture is upscaled > 2x (blurry)
  * consistency — master theme + icons present across the deck

HTML path (Marp showcase) audit via `qa2.py --html deck.html`:
  * slide count (Marp <section>), viewport meta (responsive),
  * transitions/animations present, bespoke controls present.

Outputs a JSON report + a 0-100 score. CLI:
    python3 qa2.py deck.pptx [--json report.json]
    python3 qa2.py --html deck.html
Exit code 0 when no error-level issues, else 1.
"""
from __future__ import annotations

import argparse
import json
import re
import struct
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

# Slide canvas bounds in EMU (13.333in x 7.5in) + small tolerance.
SLIDE_W = 12192000
SLIDE_H = 6858000
TOL = 9525  # ~1px tolerance for rounding

PLACEHOLDER_TOKENS = ["[", "图片缺失", "TODO", "占位", "未提供", "TBD"]
TRUNC_TOKENS = ["…", "..."]
LATIN_ONLY_FONTS = {
    "Calibri", "Arial", "Times New Roman", "Helvetica", "Cambria",
    "Georgia", "Verdana", "Tahoma", "Segoe UI", "Consolas", "Courier New",
}
DEFAULT_ACCENTS = {"5B9BD5", "4472C4", "0070C0", "1F4E79"}


def _has_cjk(s: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in s)


def _hex(c) -> str | None:
    try:
        return f"{int(c):06X}"
    except Exception:
        return None


def _rgb_tuple(hex6: str):
    h = hex6.strip("#")
    if len(h) != 6:
        return None
    try:
        return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        return None


def _luminance(rgb) -> float:
    def lin(c):
        c /= 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * lin(r) + 0.7152 * lin(g) + 0.0722 * lin(b)


def _contrast(c1, c2) -> float:
    l1, l2 = _luminance(c1), _luminance(c2)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _native_px(blob: bytes):
    if blob[:8] == b"\x89PNG\r\n\x1a\n" and len(blob) >= 24:
        w, h = struct.unpack(">II", blob[16:24])
        return w, h
    if blob[:2] == b"\xff\xd8":
        i = 2
        n = len(blob)
        while i < n - 9:
            if blob[i] != 0xFF:
                i += 1
                continue
            marker = blob[i + 1]
            if marker in (0xC0, 0xC1, 0xC2, 0xC3, 0xC5, 0xC6, 0xC7,
                          0xC9, 0xCA, 0xCB):
                h, w = struct.unpack(">HH", blob[i + 5:i + 9])
                return w, h
            seg = struct.unpack(">H", blob[i + 2:i + 4])[0]
            i += 2 + seg
    return None, None


def _pic_embeds(slide) -> list[str]:
    out = []
    for blip in slide._element.iter(qn("a:blip")):
        emb = blip.get(qn("r:embed"))
        if emb:
            out.append(emb)
    return out


def audit(path: str) -> dict:
    prs = Presentation(path)
    slides = list(prs.slides)
    total = len(slides)
    errors = []
    warnings = []
    icon_shapes = 0
    transition_slides = 0
    pictures = 0
    min_body = 999
    bounds_hits = 0
    cjk_font_hits = 0
    contrast_hits = 0
    trunc_hits = 0
    upscale_hits = 0

    for i, slide in enumerate(slides, 1):
        xml = slide._element.xml
        icon_shapes += xml.count("custGeom")
        pictures += xml.count("<p:pic>")
        if "<p:transition" in xml:
            transition_slides += 1

        # image upscale: map blip embed -> media blob -> native px
        embed_media = {}
        try:
            for rel in slide.part.rels.values():
                if "image" in rel.rel_type:
                    try:
                        embed_media[rel.rel_id] = rel.target_part.blob
                    except Exception:
                        pass
        except Exception:
            pass

        # bounds + image checks
        for sh in slide.shapes:
            try:
                l = sh.left or 0
                t = sh.top or 0
                w = sh.width or 0
                h = sh.height or 0
            except Exception:
                continue
            if (l < -TOL or t < -TOL or
                    l + w > SLIDE_W + TOL or t + h > SLIDE_H + TOL):
                bounds_hits += 1
                errors.append(
                    f"slide {i}: shape out of bounds "
                    f"(L={Emu(l).inches:.2f} T={Emu(t).inches:.2f} "
                    f"R={Emu(l + w).inches:.2f} B={Emu(t + h).inches:.2f} in)")
            # upscale check (only if we have native px + a real display size)
            if getattr(sh, "shape_type", None) is not None and "<p:pic>" in sh._element.xml:
                embeds = _pic_embeds(slide)
                for eid in embeds:
                    blob = embed_media.get(eid)
                    if not blob:
                        continue
                    nw, nh = _native_px(blob)
                    if nw and nh and w and h:
                        disp_w_px = Emu(w).inches * 96
                        disp_h_px = Emu(h).inches * 96
                        if disp_w_px > nw * 2 or disp_h_px > nh * 2:
                            upscale_hits += 1
                            warnings.append(
                                f"slide {i}: picture upscaled >2x "
                                f"(native {nw}x{nh}px, display ~{int(disp_w_px)}x{int(disp_h_px)}px)")
                    break

        # text checks
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            # shape solid fill (for contrast)
            shape_fill_hex = None
            try:
                if sh.fill.type is not None and getattr(sh.fill, "fore_color", None):
                    shape_fill_hex = _hex(sh.fill.fore_color.rgb)
            except Exception:
                shape_fill_hex = None
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    txt = run.text or ""
                    sz = run.font.size.pt if run.font.size else None
                    if txt.strip():
                        for tok in PLACEHOLDER_TOKENS:
                            if tok in txt:
                                errors.append(
                                    f"slide {i}: placeholder token '{tok}' "
                                    f"in text '{txt[:24]}'")
                                break
                        for tk in TRUNC_TOKENS:
                            if tk in txt:
                                trunc_hits += 1
                                warnings.append(
                                    f"slide {i}: text truncated ('{tk}') "
                                    f"in '{txt[:24]}'")
                                break
                        if _has_cjk(txt):
                            fn = (run.font.name or "") if run.font else ""
                            if fn in LATIN_ONLY_FONTS:
                                cjk_font_hits += 1
                                warnings.append(
                                    f"slide {i}: CJK text on Latin-only font "
                                    f"'{fn}' (may render as boxes)")
                    # contrast: run solid colour vs shape solid fill
                    if sz is not None:
                        if sz < 9:
                            errors.append(
                                f"slide {i}: font {sz}pt below 9pt floor")
                        elif 9 < sz < 11:
                            warnings.append(
                                f"slide {i}: body font {sz}pt (over-shrunk, "
                                f"target >=11)")
                        if sz >= 9:
                            min_body = min(min_body, sz)
                    run_hex = None
                    try:
                        if run.font.color and run.font.color.type is not None:
                            run_hex = _hex(run.font.color.rgb)
                    except Exception:
                        run_hex = None
                    if run_hex and shape_fill_hex:
                        c1 = _rgb_tuple(run_hex)
                        c2 = _rgb_tuple(shape_fill_hex)
                        if c1 and c2 and _contrast(c1, c2) < 3.0:
                            contrast_hits += 1
                            warnings.append(
                                f"slide {i}: low contrast text/fill "
                                f"(ratio {_contrast(c1, c2):.2f} < 3.0)")

    # brand theme check (theme-agnostic)
    brand = False
    notes_slides = 0
    try:
        z = zipfile.ZipFile(path)
        th = [n for n in z.namelist() if re.search(r"theme1?\.xml$", n)]
        if th:
            x = z.read(th[0]).decode("utf-8", "ignore")
            m = re.search(r"<a:accent1>\s*<a:srgbClr val=\"([0-9A-Fa-f]{6})\"", x)
            if m and m.group(1).upper() not in DEFAULT_ACCENTS:
                brand = True
        notes_slides = len([n for n in z.namelist()
                            if re.match(r"ppt/notesSlides/notesSlide\d+\.xml$", n)])
    except Exception:
        pass

    # score
    score = 100
    score -= len(errors) * 8
    score -= len(warnings) * 2
    if transition_slides != total:
        score -= 5
    if not brand:
        score -= 5
    if cjk_font_hits:
        score -= 4
    if contrast_hits:
        score -= 3
    if trunc_hits:
        score -= 2
    if upscale_hits:
        score -= 2
    score = max(0, min(100, score))

    return {
        "file": str(Path(path).resolve()),
        "slides": total,
        "slides_clean": total - (1 if errors else 0),
        "errors": errors,
        "warnings": warnings,
        "bounds_hits": bounds_hits,
        "icon_shapes": icon_shapes,
        "pictures": pictures,
        "transitions": transition_slides,
        "notes_slides": notes_slides,
        "min_body_pt": min_body if min_body != 999 else None,
        "cjk_font_hits": cjk_font_hits,
        "contrast_hits": contrast_hits,
        "trunc_hits": trunc_hits,
        "upscale_hits": upscale_hits,
        "brand_theme": brand,
        "score": score,
    }


def audit_html(path: str) -> dict:
    """Audit a Marp/HTML showcase deck for QA dimensions that matter there."""
    html = Path(path).read_text(encoding="utf-8", errors="ignore")
    sections = html.count("<section")
    has_viewport = 'name="viewport"' in html or 'name=\'viewport\'' in html
    has_anim = ("animation" in html) or ("transition" in html)
    has_bespoke = "bespoke" in html
    # crude mobile responsiveness: a CSS rule with max-width or vw units
    responsive = bool(re.search(r"max-width\s*:|vw|@media", html))
    errors = []
    warnings = []
    if sections == 0:
        errors.append("no <section> slides found (Marp deck empty?)")
    if not has_viewport:
        warnings.append("missing viewport meta (mobile may mis-scale)")
    if not has_anim:
        warnings.append("no transition/animation CSS (less 'showcase' feel)")
    score = 100 - len(errors) * 10 - len(warnings) * 3
    score = max(0, min(100, score))
    return {
        "file": str(Path(path).resolve()),
        "slides": sections,
        "viewport_meta": has_viewport,
        "animations": has_anim,
        "bespoke_controls": has_bespoke,
        "responsive": responsive,
        "errors": errors,
        "warnings": warnings,
        "score": score,
    }


def main():
    ap = argparse.ArgumentParser(description="PPT Pro Studio QA 2.1 audit")
    ap.add_argument("pptx", nargs="?", help="path to .pptx (or .html with --html)")
    ap.add_argument("--html", action="store_true", help="audit an HTML/Marp deck")
    ap.add_argument("--json", default=None, help="write JSON report to this path")
    args = ap.parse_args()
    if not args.pptx:
        sys.stderr.write("usage: qa2.py [--html] <file>\n")
        sys.exit(2)
    if not Path(args.pptx).exists():
        sys.stderr.write(f"file not found: {args.pptx}\n")
        sys.exit(2)

    if args.html:
        rep = audit_html(args.pptx)
        if args.json:
            Path(args.json).write_text(json.dumps(rep, ensure_ascii=False, indent=2),
                                       encoding="utf-8")
        print("=" * 56)
        print(f"QA 2.1 (HTML)  ·  {rep['file']}")
        print(f"slides={rep['slides']}  score={rep['score']}/100")
        print(f"viewport={rep['viewport_meta']}  animations={rep['animations']}  "
              f"responsive={rep['responsive']}")
        if rep["errors"]:
            print(f"\nERRORS ({len(rep['errors'])}):")
            for e in rep["errors"]:
                print("  ✗", e)
        if rep["warnings"]:
            print(f"\nWARNINGS ({len(rep['warnings'])}):")
            for w in rep["warnings"]:
                print("  !", w)
        if not rep["errors"] and not rep["warnings"]:
            print("\n✓ no HTML issues found")
        print("=" * 56)
        sys.exit(1 if rep["errors"] else 0)

    rep = audit(args.pptx)
    if args.json:
        Path(args.json).write_text(json.dumps(rep, ensure_ascii=False, indent=2),
                                   encoding="utf-8")
    print("=" * 56)
    print(f"QA 2.1  ·  {rep['file']}")
    print(f"slides={rep['slides']}  score={rep['score']}/100")
    print(f"clean={rep['slides_clean']}  bounds_hits={rep['bounds_hits']}  "
          f"transitions={rep['transitions']}/{rep['slides']}")
    print(f"icon/vector shapes={rep['icon_shapes']}  pictures={rep['pictures']}  "
          f"notes_slides={rep['notes_slides']}")
    print(f"min_body={rep['min_body_pt']}pt  brand_theme={rep['brand_theme']}")
    print(f"cjk_font_hits={rep['cjk_font_hits']}  contrast_hits={rep['contrast_hits']}  "
          f"trunc_hits={rep['trunc_hits']}  upscale_hits={rep['upscale_hits']}")
    if rep["errors"]:
        print(f"\nERRORS ({len(rep['errors'])}):")
        for e in rep["errors"]:
            print("  ✗", e)
    if rep["warnings"]:
        print(f"\nWARNINGS ({len(rep['warnings'])}):")
        for w in rep["warnings"]:
            print("  !", w)
    if not rep["errors"] and not rep["warnings"]:
        print("\n✓ no layout issues found")
    print("=" * 56)
    sys.exit(1 if rep["errors"] else 0)


if __name__ == "__main__":
    main()
