#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ppt_studio_generate.py — PPT Pro Studio core renderer.

Turns a structured JSON *brief* into a commercial-grade, fully editable .pptx.
Designed to be called by:
  * the ppt-pro-studio SKILL (LLM agents following the workflow)
  * the ppt-studio-mcp server (any MCP-capable LLM)
  * directly from the CLI

No network, no external services, no watermarks. MIT-0.
"""
from __future__ import annotations

import argparse
import json
import os
import subprocess
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn

try:
    from icons import add_icon_to_slide, DEFAULT_ICON
except ImportError:  # allow running with scripts/ as cwd
    import os
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from icons import add_icon_to_slide, DEFAULT_ICON

try:
    from theme_market import load_themes, select_theme
except ImportError:  # allow running with scripts/ as cwd
    import os
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from theme_market import load_themes, select_theme

# ---------------------------------------------------------------------------
# Design system — 5 commercial palettes (+ the dark-tech reference palette)
# Every palette: bg, surface, primary, secondary, text, muted, accent, line
# ---------------------------------------------------------------------------

# 8pt grid + type scale tokens (design-system 2.0). All spacing in inches
# is kept on an 8px-multiple rhythm; font sizes follow a strict scale so the
# hierarchy reads consistently across every slide.
GRID = 0.05          # base grid unit (in) ~ 4px
SPACING = {          # named spacing tokens (in)
    "xs": 0.1, "sm": 0.2, "md": 0.3, "lg": 0.5, "xl": 0.8, "xxl": 1.2,
}
TYPE = {             # type scale (pt)
    "caption": 9, "body": 15, "body_sm": 12, "title": 28,
    "subtitle": 14, "display": 44, "section": 34, "hero": 46,
}
FONT_FLOOR = 11      # never shrink body text below this (readability floor)


def _char_w_pt(ch: str, sz: float) -> float:
    """Approx glyph advance in points: CJK ~full em, ASCII ~0.55em."""
    return sz if ord(ch) > 0x2E80 else sz * 0.55


def _estimate_lines(text: str, sz: float, w_in: float) -> int:
    line_w = w_in * 72.0
    lines = 0
    for para in text.split("\n"):
        if not para:
            lines += 1
            continue
        w = 0.0
        l = 1
        for ch in para:
            cw = _char_w_pt(ch, sz)
            if w + cw > line_w and w > 0:
                l += 1
                w = cw
            else:
                w += cw
        lines += l
    return max(1, lines)


def _fit_text(text: str, w_in: float, h_in: float, size: float, floor=FONT_FLOOR):
    """Pick the largest font >= floor that fits the box; if even the floor
    overflows, truncate the text (adding an ellipsis) instead of shrinking
    into an unreadable size."""
    sz = float(size)
    while True:
        need = _estimate_lines(text, sz, w_in) * sz * 1.25
        if need <= h_in * 72.0 or sz <= floor:
            break
        sz -= 1.0
    sz = max(sz, float(floor))
    if _estimate_lines(text, sz, w_in) * sz * 1.25 > h_in * 72.0:
        while len(text) > 1 and \
                _estimate_lines(text, sz, w_in) * sz * 1.25 > h_in * 72.0:
            text = text[:-1]
        text = text.rstrip() + "\u2026"
    return sz, text

PALETTES = {
    "tech_dark": {  # 深色科技风 (from the commercial delivery reference pack)
        "bg": "0D1117", "surface": "151D2E", "primary": "D4A060", "secondary": "58A6FF",
        "text": "FFFFFF", "muted": "8B949E", "accent": "3FB950", "line": "30363D",
        "font": "Microsoft YaHei",
    },
    "business_blue": {
        "bg": "FFFFFF", "surface": "F2F6FC", "primary": "1F4E79", "secondary": "2E75B6",
        "text": "1A1A1A", "muted": "5A5A5A", "accent": "C55A11", "line": "D6E0F0",
        "font": "Microsoft YaHei",
    },
    "creative_purple": {
        "bg": "1B1033", "surface": "2A1B4A", "primary": "C77DFF", "secondary": "7B2FBE",
        "text": "F5EEFF", "muted": "B39DCE", "accent": "FF8FB1", "line": "3D2A63",
        "font": "Microsoft YaHei",
    },
    "academic_white": {
        "bg": "FFFFFF", "surface": "FAFAFA", "primary": "202020", "secondary": "006633",
        "text": "1A1A1A", "muted": "666666", "accent": "B00020", "line": "DDDDDD",
        "font": "Microsoft YaHei",
    },
    "minimal_gray": {
        "bg": "FAFAFA", "surface": "F0F0F0", "primary": "222222", "secondary": "666666",
        "text": "1A1A1A", "muted": "999999", "accent": "007ACC", "line": "E0E0E0",
        "font": "Microsoft YaHei",
    },
}

STYLE_LABELS = {
    "tech_dark": "深色科技风 (Tech Dark)",
    "business_blue": "商务蓝 (Business Blue)",
    "creative_purple": "创意紫 (Creative Purple)",
    "academic_white": "学术白 (Academic White)",
    "minimal_gray": "简约灰 (Minimal Gray)",
}

# Theme market: on-disk themes/*.json merged OVER the built-in PALETTES above,
# so users can add custom themes without editing engine code. The brief picks a
# theme via `theme` (or the `style` alias).
THEMES = load_themes(builtin=PALETTES)
STYLE_LABELS = {k: v.get("label", k) for k, v in THEMES.items()}

FONT = "Microsoft YaHei"  # cross-platform: falls back to default on non-Windows
EMU_IN = 914400
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _c(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


class Studio:
    def __init__(self, palette_name: str | None = None):
        self.p = THEMES.get(palette_name, THEMES["tech_dark"])
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]

    # -- master / theme -----------------------------------------------------
    def _apply_master_theme(self, out_path):
        """Write brand colors + font into the pptx theme (theme1.xml) so they
        are editable via the PowerPoint/WPS 'Colors' and 'Fonts' panels (one
        place to restyle the whole deck). Patches the saved .pptx zip directly
        to stay robust across python-pptx builds. Standard OOXML, WPS-safe."""
        try:
            import zipfile
            import shutil
            from lxml import etree
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            pal = self.p
            order = [("dk1", pal["bg"]), ("lt1", pal["text"]),
                     ("dk2", pal["secondary"]), ("lt2", pal["surface"]),
                     ("accent1", pal["primary"]), ("accent2", pal["secondary"]),
                     ("accent3", pal["accent"]), ("accent4", pal["muted"]),
                     ("accent5", pal["line"]), ("accent6", pal["primary"])]
            tmp = out_path + ".tmp"
            src = zipfile.ZipFile(out_path, "r")
            theme_name = [n for n in src.namelist()
                          if __import__("re").search(r"theme1?\.xml$", n)]
            if not theme_name:
                src.close()
                return
            theme_name = theme_name[0]
            root = etree.fromstring(src.read(theme_name))
            te = root.find(f"{{{ns}}}themeElements")
            if te is None:
                src.close()
                return
            cs = te.find(f"{{{ns}}}clrScheme")
            if cs is not None:
                for c in list(cs):
                    cs.remove(c)
                for tag, hexv in order:
                    e = etree.SubElement(cs, f"{{{ns}}}{tag}")
                    c = etree.SubElement(e, f"{{{ns}}}srgbClr")
                    c.set("val", hexv)
            fs = te.find(f"{{{ns}}}fontScheme")
            if fs is not None:
                for which in ("majorFont", "minorFont"):
                    mf = fs.find(f"{{{ns}}}{which}")
                    if mf is None:
                        continue
                    latin = mf.find(f"{{{ns}}}latin")
                    if latin is None:
                        latin = etree.SubElement(mf, f"{{{ns}}}latin")
                    latin.set("typeface", pal["font"])
            new_theme = etree.tostring(root, xml_declaration=True,
                                       encoding="UTF-8", standalone=True)
            with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zo:
                for item in src.infolist():
                    if item.filename == theme_name:
                        zo.writestr(item, new_theme)
                    else:
                        zo.writestr(item, src.read(item.filename))
            src.close()
            shutil.move(tmp, out_path)
        except Exception as e:  # theme is best-effort, never fatal
            sys.stderr.write("warning: master theme skipped: %s\n" % e)

    def _transparent_chart(self, chart):
        """Make the chart plot area transparent so the brand slide background
        shows through (no white box on dark themes). Standard OOXML patch,
        WPS-safe. Best-effort, never fatal."""
        try:
            from pptx.oxml.ns import qn
            cs = chart._chartSpace
            el = cs.find(qn("c:plotArea"))
            if el is None:
                return
            spPr = el.find(qn("c:spPr"))
            if spPr is None:
                spPr = el.makeelement(qn("c:spPr"), {})
                el.insert(0, spPr)
            for child in list(spPr):
                if child.tag in (qn("a:solidFill"), qn("a:gradFill"),
                                 qn("a:blipFill"), qn("a:pattFill"),
                                 qn("a:grpFill")):
                    spPr.remove(child)
            if spPr.find(qn("a:noFill")) is None:
                spPr.append(spPr.makeelement(qn("a:noFill"), {}))
        except Exception:
            pass

    def _icon(self, slide, name, x, y, size, color):
        try:
            add_icon_to_slide(slide, name, x, y, size, color, stroke_w_pt=1.5)
        except Exception:
            pass

    # -- low level helpers ---------------------------------------------------
    def _slide(self):
        return self.prs.slides.add_slide(self._blank)

    def _bg(self, slide, color=None):
        color = color or self.p["bg"]
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _c(color)

    def _txt(self, slide, x, y, w, h, text, size=18, color=None, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font=None,
             floor=FONT_FLOOR):
        # Enforce readability floor + line-length control: shrink only down to
        # `floor`, then truncate with an ellipsis instead of going tiny.
        size, text = _fit_text(text, w, h, size, floor=floor)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        lines = text.split("\n")
        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align
            run = para.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = font or self.p["font"]
            run.font.color.rgb = _c(color or self.p["text"])
        return tb

    def _rect(self, slide, x, y, w, h, fill, line=None, line_w=1):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                     Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _c(fill)
        if line:
            shp.line.color.rgb = _c(line)
            shp.line.width = Pt(line_w)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _footer(self, slide, idx, footer, page_numbers):
        if footer:
            self._txt(slide, 0.5, 7.05, 9, 0.35, footer, size=9,
                      color=self.p["muted"], floor=9)
        if page_numbers:
            self._txt(slide, 12.3, 7.05, 0.8, 0.35, str(idx), size=9,
                      color=self.p["muted"], align=PP_ALIGN.RIGHT, floor=9)

    def _title_bar(self, slide, title, subtitle=None, num=None, stype="content",
                   icon=None):
        name = icon or DEFAULT_ICON.get(stype, "list")
        self._icon(slide, name, 0.55, 0.42, 0.5, self.p["primary"])
        self._txt(slide, 1.2, 0.4, 11.0, 0.7, title, size=28, bold=True,
                  color=self.p["primary"])
        self._rect(slide, 1.2, 1.12, 2.2, 0.06, self.p["primary"])
        if subtitle:
            self._txt(slide, 1.2, 1.22, 11.0, 0.45, subtitle, size=14,
                      color=self.p["secondary"])
        if num:
            self._txt(slide, 11.8, 0.4, 1.0, 0.7, num, size=22, bold=True,
                      color=self.p["muted"], align=PP_ALIGN.RIGHT)

    # -- slide renderers -----------------------------------------------------
    def cover(self, s, data):
        self._bg(s)
        variant = data.get("variant", "centered")
        title = data.get("title", "")
        subtitle = data.get("subtitle", "")
        badge = data.get("badge", "")
        icon = data.get("icon", "rocket")
        if variant == "left":
            self._icon(s, icon, 0.9, 2.4, 0.9, self.p["primary"])
            self._txt(s, 0.9, 3.5, 11, 1.4, title, size=44, bold=True,
                      color=self.p["primary"])
            if subtitle:
                self._txt(s, 0.9, 5.0, 10, 0.8, subtitle, size=20,
                          color=self.p["secondary"])
        else:  # centered
            self._icon(s, icon, 6.27, 1.5, 0.8, self.p["primary"])
            self._txt(s, 1, 2.6, 11.3, 1.5, title, size=46, bold=True,
                      color=self.p["primary"], align=PP_ALIGN.CENTER)
            if subtitle:
                self._txt(s, 1, 4.2, 11.3, 0.8, subtitle, size=20,
                          color=self.p["secondary"], align=PP_ALIGN.CENTER)
        if badge:
            self._txt(s, 1, 5.5, 11.3, 0.5, badge, size=13,
                      color=self.p["muted"], align=PP_ALIGN.CENTER)

    def section(self, s, data):
        self._bg(s)
        title = data.get("title", "")
        num = data.get("index") or data.get("num")
        title_text = title
        if not num:
            sp = title.split(" ", 1)
            if sp[0].isdigit():
                num = sp[0]
                title_text = sp[1] if len(sp) > 1 else title
        # Band + accent edge
        self._rect(s, 0, 2.7, 13.333, 2.0, self.p["surface"])
        self._rect(s, 0, 2.7, 0.18, 2.0, self.p["primary"])
        # Vertical divider separating the big index from the title
        self._rect(s, 3.4, 2.9, 0.02, 1.6, self.p["line"])
        # Big index number (narrative rhythm)
        if num:
            self._txt(s, 0.6, 2.7, 2.6, 2.0, num, size=96, bold=True,
                      color=self.p["accent"], align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)
        # Geometric ring decoration (right)
        ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.3), Inches(2.95),
                                  Inches(1.3), Inches(1.3))
        ring.fill.background()
        ring.line.color.rgb = _c(self.p["primary"])
        ring.line.width = Pt(1.5)
        ring.shadow.inherit = False
        # Title + subtitle (right of the divider)
        tx = 3.7 if num else 1.7
        self._txt(s, tx, 2.95, 7.4, 1.0, title_text, size=34, bold=True,
                  color=self.p["primary"], anchor=MSO_ANCHOR.MIDDLE)
        if data.get("subtitle"):
            self._txt(s, tx, 4.35, 7.4, 0.5, data["subtitle"], size=14,
                      color=self.p["muted"])

    def agenda(self, s, data):
        self._bg(s)
        self._title_bar(s, data.get("title", "目录"), data.get("subtitle"),
                       stype="agenda")
        items = data.get("items", [])
        top, bottom = 1.7, 6.9
        gap = min(0.78, (bottom - top) / max(1, len(items)))
        for i, it in enumerate(items):
            y = top + i * gap
            self._txt(s, 0.9, y, 0.7, 0.6, f"{i+1:02d}", size=20, bold=True,
                      color=self.p["secondary"])
            self._txt(s, 1.7, y, 10.5, 0.6, it, size=16, color=self.p["text"])
            self._rect(s, 0.9, y + min(gap, 0.78) - 0.12, 11.4, 0.015, self.p["line"])

    def content(self, s, data):
        self._bg(s)
        self._title_bar(s, data.get("title", ""), data.get("subtitle"),
                        data.get("num"), stype="content")
        items = data.get("items", [])
        columns = data.get("columns", 1)
        if columns == 2 and len(items) > 3:
            half = (len(items) + 1) // 2
            self._bullets(s, items[:half], 0.9, 1.7, 5.8, 0.62)
            self._bullets(s, items[half:], 6.9, 1.7, 5.8, 0.62)
        else:
            self._bullets(s, items, 0.9, 1.7, 11.5, 0.62)

    def _bullets(self, s, items, x, y, w, gap=0.62, max_y=7.0):
        if items:
            # Shrink the per-item gap so many bullets never overflow the slide.
            avail = max_y - y
            gap = min(gap, avail / len(items))
        for i, it in enumerate(items):
            yy = y + i * gap
            # bullet dot
            self._rect(s, x, yy + 0.12, 0.12, 0.12, self.p["secondary"])
            # support nested "detail" via "|"
            if "|" in it:
                head, detail = it.split("|", 1)
            else:
                head, detail = it, ""
            self._txt(s, x + 0.3, yy, w - 0.3, 0.55, head, size=15,
                      color=self.p["text"], bold=True)
            if detail:
                self._txt(s, x + 0.3, yy + 0.42, w - 0.3, 0.5, detail, size=12,
                          color=self.p["muted"])

    def two_column(self, s, data):
        self._bg(s)
        self._title_bar(s, data.get("title", ""), data.get("subtitle"),
                       stype="two_column")
        self._rect(s, 0.6, 1.7, 5.9, 4.8, self.p["surface"])
        self._rect(s, 6.8, 1.7, 5.9, 4.8, self.p["surface"])
        self._txt(s, 0.9, 1.9, 5.3, 0.5, data.get("left_title", "左栏"),
                  size=16, bold=True, color=self.p["primary"])
        self._txt(s, 7.1, 1.9, 5.3, 0.5, data.get("right_title", "右栏"),
                  size=16, bold=True, color=self.p["primary"])
        self._bullets(s, data.get("left", []), 0.9, 2.5, 5.3, 0.6)
        self._bullets(s, data.get("right", []), 7.1, 2.5, 5.3, 0.6)

    def table(self, s, data):
        self._bg(s)
        self._title_bar(s, data.get("title", ""), data.get("subtitle"),
                       stype="table")
        headers = data.get("headers", [])
        rows = data.get("rows", [])
        nrows = len(rows) + 1
        ncols = max(len(headers), *(len(r) for r in rows) if rows else [1])
        gtab = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(1.7),
                                  Inches(12.1), Inches(4.8)).table
        gtab.columns  # ensure
        for j, h in enumerate(headers):
            cell = gtab.cell(0, j)
            cell.text = str(h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _c(self.p["primary"])
            self._cell_font(cell, size=13, bold=True, color="FFFFFF")
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                cell = gtab.cell(i, j)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _c(self.p["surface"] if i % 2 else self.p["bg"])
                self._cell_font(cell, size=12, color=self.p["text"])

    def _cell_font(self, cell, size=12, bold=False, color="FFFFFF"):
        tf = cell.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.name = self.p["font"]
                run.font.color.rgb = _c(color)

    def chart(self, s, data):
        self._bg(s)
        self._title_bar(s, data.get("title", ""), data.get("subtitle"),
                       stype="chart")
        ctype = data.get("chart_type", "bar").lower()
        mapping = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                   "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                   "line": XL_CHART_TYPE.LINE_MARKERS,
                   "pie": XL_CHART_TYPE.PIE}
        chart_data = CategoryChartData()
        chart_data.categories = data.get("categories", [])
        for ser in data.get("series", []):
            chart_data.add_series(ser.get("name", "Series"), ser.get("values", []))
        gf = s.shapes.add_chart(mapping.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED),
                                Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6),
                                chart_data)
        chart = gf.chart
        # Transparent chart/plot area so the brand background shows through.
        try:
            chart.fill.background()
            chart.plot_area.fill.background()
        except Exception:
            pass
        self._transparent_chart(chart)
        # Legend: only when >1 series, placed at the bottom, theme-colored.
        chart.has_legend = len(data.get("series", [])) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
            chart.legend.font.name = self.p["font"]
            chart.legend.font.color.rgb = _c(self.p["text"])
        chart.font.name = self.p["font"]
        chart.font.size = Pt(11)
        chart.font.color.rgb = _c(self.p["text"])
        # Axes: gridlines + theme-colored tick labels for readable contrast.
        try:
            val = chart.value_axis
            val.has_major_gridlines = True
            val.major_gridlines.format.line.color.rgb = _c(self.p["line"])
            val.major_gridlines.format.line.width = Pt(0.5)
            val.tick_labels.font.size = Pt(10)
            val.tick_labels.font.color.rgb = _c(self.p["muted"])
            val.has_title = False
            cat = chart.category_axis
            cat.tick_labels.font.size = Pt(11)
            cat.tick_labels.font.color.rgb = _c(self.p["text"])
        except Exception:
            pass
        # Series colors (distinct), data labels, and single-point emphasis.
        series_colors = [self.p["primary"], self.p["secondary"], self.p["accent"]]
        emphasis = data.get("emphasis", [])
        if isinstance(emphasis, int):
            emphasis = [emphasis]
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            plot.data_labels.number_format = "0"
            plot.data_labels.number_format_is_linked = False
            plot.data_labels.font.size = Pt(10)
            plot.data_labels.font.name = self.p["font"]
            plot.data_labels.font.color.rgb = _c(self.p["text"])
            if ctype == "pie":
                plot.data_labels.show_percentage = True
                plot.data_labels.show_category_name = False
                plot.data_labels.show_value = False
            elif ctype != "line":
                plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            for si, ser in enumerate(chart.series):
                col = series_colors[si % len(series_colors)]
                if ctype == "line":
                    try:
                        ser.format.line.color.rgb = _c(col)
                        ser.format.line.width = Pt(2.25)
                    except Exception:
                        pass
                else:
                    try:
                        ser.format.fill.solid()
                        ser.format.fill.fore_color.rgb = _c(col)
                    except Exception:
                        pass
                try:
                    for pi in emphasis:
                        pt = ser.points[pi]
                        pt.format.fill.solid()
                        pt.format.fill.fore_color.rgb = _c(self.p["accent"])
                except Exception:
                    pass
        except Exception:
            pass

    def timeline(self, s, data):
        self._bg(s)
        self._title_bar(s, data.get("title", ""), data.get("subtitle"),
                       stype="timeline")
        ms = data.get("milestones", [])
        if not ms:
            return
        y = 4.0
        self._rect(s, 1.0, y, 11.3, 0.04, self.p["line"])
        n = len(ms)
        step = 11.3 / n
        for i, m in enumerate(ms):
            cx = 1.0 + step * (i + 0.5)
            self._rect(s, cx - 0.09, y - 0.07, 0.18, 0.18, self.p["primary"])
            self._txt(s, cx - 1.2, y - 1.1, 2.4, 0.9, m.get("label", ""),
                      size=13, bold=True, color=self.p["primary"],
                      align=PP_ALIGN.CENTER)
            self._txt(s, cx - 1.2, y + 0.15, 2.4, 1.0, m.get("desc", ""),
                      size=11, color=self.p["muted"], align=PP_ALIGN.CENTER)

    def quote(self, s, data):
        self._bg(s)
        self._icon(s, data.get("icon", "quote"), 0.7, 1.6, 0.9, self.p["primary"])
        self._txt(s, 1.2, 1.7, 1.5, 1.5, "\u201C", size=120, bold=True,
                  color=self.p["primary"])
        self._txt(s, 1.5, 2.3, 10.5, 2.5, data.get("quote", ""), size=26,
                  italic=True, color=self.p["text"], anchor=MSO_ANCHOR.MIDDLE)
        if data.get("attribution"):
            self._txt(s, 1.5, 5.2, 10.5, 0.6, f"— {data['attribution']}",
                      size=15, color=self.p["secondary"])

    def image(self, s, data):
        self._bg(s)
        self._title_bar(s, data.get("title", ""), data.get("subtitle"),
                       stype="image")
        path = data.get("image_path", "")
        x, y, w, h = 1.2, 1.8, 10.9, 4.6
        if path and os.path.exists(path):
            s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
        else:
            self._rect(s, x, y, w, h, self.p["surface"], line=self.p["line"])
            self._txt(s, x, y + h / 2 - 0.3, w, 0.6,
                      "[ 图片缺失: %s ]" % path, size=14,
                      color=self.p["muted"], align=PP_ALIGN.CENTER)

    def media(self, s, data):
        """图文混排 (image + text mixed layout). Image sits left/right, the
        text block (heading + bullets, or free text) fills the other side.
        Image is contain-fitted via its natural aspect ratio (PIL best-effort)
        and centered in its column; missing image -> branded placeholder."""
        self._bg(s)
        self._title_bar(s, data.get("title", ""), data.get("subtitle"),
                       stype="media")
        img = data.get("image") or data.get("image_path") or ""
        pos = (data.get("image_position") or "left").lower()
        img_x, img_y, img_w, img_h = 0.6, 1.8, 5.7, 4.6
        txt_x = 6.7
        if pos == "right":
            img_x, txt_x = 6.9, 0.6
        if img and os.path.exists(img):
            try:
                from PIL import Image
                with Image.open(img) as im:
                    iw, ih = im.size
                box_ar, nat_ar = img_w / img_h, (iw / ih) if ih else 1.0
                if nat_ar > box_ar:
                    w, h = img_w, img_w / nat_ar
                else:
                    h, w = img_h, img_h * nat_ar
                left = img_x + (img_w - w) / 2.0
                top = img_y + (img_h - h) / 2.0
                s.shapes.add_picture(img, Inches(left), Inches(top),
                                     Inches(w), Inches(h))
            except Exception:
                self._rect(s, img_x, img_y, img_w, img_h, self.p["surface"],
                           line=self.p["line"])
                self._txt(s, img_x, img_y + img_h / 2 - 0.3, img_w, 0.6,
                          "[ 图片缺失 ]", size=14, color=self.p["muted"],
                          align=PP_ALIGN.CENTER)
        else:
            self._rect(s, img_x, img_y, img_w, img_h, self.p["surface"],
                       line=self.p["line"])
            self._txt(s, img_x, img_y + img_h / 2 - 0.3, img_w, 0.6,
                      "[ 图片缺失 ]", size=14, color=self.p["muted"],
                      align=PP_ALIGN.CENTER)
        if data.get("caption"):
            self._txt(s, img_x, img_y + img_h + 0.05, img_w, 0.4,
                      data["caption"], size=11, color=self.p["muted"],
                      align=PP_ALIGN.CENTER)
        ty = 1.9
        if data.get("heading"):
            self._txt(s, txt_x, ty, 6.0, 0.6, data["heading"], size=18,
                      bold=True, color=self.p["primary"])
            ty += 0.75
        if data.get("text"):
            self._txt(s, txt_x, ty, 6.0, (1.9 + 4.6) - ty, data["text"],
                      size=14, color=self.p["text"])
        else:
            self._bullets(s, data.get("items", []), txt_x, ty, 6.0, 0.6)

    def summary(self, s, data):
        self._bg(s)
        self._title_bar(s, data.get("title", "总结"), stype="summary")
        self._bullets(s, data.get("points", []), 0.9, 1.7, 11.5, 0.66)
        if data.get("conclusion"):
            self._rect(s, 0.9, 6.0, 11.5, 0.9, self.p["primary"])
            self._txt(s, 1.1, 6.0, 11.1, 0.9, data["conclusion"], size=15,
                      bold=True, color="FFFFFF", anchor=MSO_ANCHOR.MIDDLE)

    def bullets(self, s, data):
        self._bg(s)
        self._title_bar(s, data.get("title", ""), data.get("subtitle"),
                       stype="bullets")
        self._bullets(s, data.get("items", []), 0.9, 1.7, 11.5, 0.66)

    def contact(self, s, data):
        self._bg(s)
        self._icon(s, data.get("icon", "mail"), 6.27, 1.4, 0.8,
                   self.p["primary"])
        self._txt(s, 1, 2.4, 11.3, 1.0, data.get("title", "联系我们"),
                  size=36, bold=True, color=self.p["primary"],
                  align=PP_ALIGN.CENTER)
        self._txt(s, 1, 3.6, 11.3, 1.5, data.get("info", ""), size=16,
                  color=self.p["text"], align=PP_ALIGN.CENTER)

    # -- dispatch ------------------------------------------------------------
    RENDERERS = {
        "cover": cover, "section": section, "agenda": agenda,
        "content": content, "two_column": two_column, "table": table,
        "chart": chart, "timeline": timeline, "quote": quote,
        "image": image, "media": media, "summary": summary, "bullets": bullets,
        "contact": contact,
    }

    def build(self, brief: dict, out_path: str, no_transition: bool = False,
              seed=None, duration: float = 0.5):
        self.p = select_theme(brief, THEMES)
        style = brief.get("theme") or brief.get("style") or "tech_dark"
        slides = brief.get("slides", [])
        total = len(slides)
        for idx, sl in enumerate(slides, 1):
            stype = sl.get("type", "content")
            renderer = self.RENDERERS.get(stype, self.content)
            s = self._slide()
            renderer(self, s, sl)
            # speaker notes (演讲者备注) — best effort, never fatal
            notes = sl.get("notes")
            if notes:
                try:
                    s.notes_slide.notes_text_frame.text = notes
                except Exception:
                    pass
            # footer + page numbers on non-cover/section pages
            if stype not in ("cover", "section"):
                self._footer(s, idx, brief.get("footer", ""),
                             brief.get("page_numbers", True))
        os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
        self.prs.save(out_path)

        # Best-effort: brand the slide-master theme (colors + fonts editable
        # via the PPT/WPS panels). Patches the saved zip; never fatal.
        try:
            self._apply_master_theme(out_path)
        except Exception as e:  # pragma: no cover - defensive
            sys.stderr.write("warning: master theme skipped: %s\n" % e)

        # Best-effort: inject random per-slide page transitions (翻页随机动画).
        # Uses ppt-master's vetted OOXML transition core via add_transitions.py.
        transitions = False
        if not no_transition:
            try:
                here = Path(__file__).resolve().parent
                at = here / "add_transitions.py"
                if at.exists():
                    targs = [sys.executable, str(at), out_path,
                             "--duration", str(duration)]
                    if seed is not None:
                        targs += ["--seed", str(seed)]
                    r = subprocess.run(
                        targs,
                        capture_output=True, text=True, encoding="utf-8",
                    )
                    transitions = r.returncode == 0
                    if not transitions:
                        sys.stderr.write("warning: transitions skipped: "
                                         + (r.stderr or r.stdout) + "\n")
            except Exception as e:  # transition is best-effort, never fatal
                sys.stderr.write("warning: transitions skipped: %s\n" % e)
        return out_path, transitions


def main():
    ap = argparse.ArgumentParser(description="PPT Pro Studio renderer")
    ap.add_argument("brief", help="path to brief JSON")
    ap.add_argument("--out", default="output.pptx", help="output .pptx path")
    ap.add_argument("--no-transition", action="store_true",
                    help="skip random page-transition injection")
    ap.add_argument("--seed", type=int, default=None,
                    help="random seed for transition effects")
    ap.add_argument("--transition-duration", type=float, default=0.5)
    args = ap.parse_args()
    with open(args.brief, "r", encoding="utf-8") as f:
        brief = json.load(f)
    studio = Studio(brief.get("theme") or brief.get("style"))
    out, transitions = studio.build(
        brief, args.out, no_transition=args.no_transition,
        seed=args.seed, duration=args.transition_duration)
    print(json.dumps({"ok": True, "file": os.path.abspath(out),
                      "slides": len(brief.get("slides", [])),
                      "transitions": transitions,
                      "engine": "python-pptx (deterministic fallback)",
                      "theme": brief.get("theme") or brief.get("style", "tech_dark")},
                     ensure_ascii=False))


if __name__ == "__main__":
    main()
